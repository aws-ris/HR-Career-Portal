import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Elegant executive color palette
    navy = RGBColor(15, 23, 42)         # #0F172A (Deep Slate/Navy)
    indigo = RGBColor(79, 70, 229)      # #4F46E5 (Indigo)
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(248, 250, 252) # #F8FAFC (Soft Background)
    text_dark = RGBColor(30, 41, 59)     # #1E293B (Slate Dark)
    text_muted = RGBColor(100, 116, 139) # #64748B (Slate Muted)
    red = RGBColor(220, 38, 38)          # #DC2626 (Executive Red)
    green = RGBColor(22, 163, 74)        # #16A34A (Executive Green)

    # Background helper
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Title helper
    def add_slide_title(slide, text, color=navy):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color

    # Helper to add left text and right image
    def add_text_and_image(slide, title, bullets, image_filename):
        add_slide_title(slide, title)
        
        # Left Text Box (widened slightly to prevent premature wrapping)
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.4), Inches(5.0))
        tf = t_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        for idx, item in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            if isinstance(item, tuple):
                p.text = f"• {item[0]}: "
                p.font.bold = True
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = navy
                
                run = p.add_run()
                run.text = item[1]
                run.font.bold = False
                run.font.name = 'Arial'
                run.font.size = Pt(13)
                run.font.color.rgb = text_dark
            else:
                p.text = f"• {item}"
                p.font.name = 'Arial'
                p.font.size = Pt(13)
                p.font.color.rgb = text_dark
            p.space_after = Pt(12)

        # Right Image
        brain_dir = "C:/Users/Viraal/.gemini/antigravity-ide/brain/b52771bd-afd4-41e0-baec-88c88ae23617"
        image_path = os.path.join(brain_dir, image_filename)
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(6.5), Inches(1.5), Inches(6.0), Inches(4.8))
        else:
            print(f"Warning: Image {image_filename} not found at {image_path}")
            fallback_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6.0), Inches(4.8))
            ftf = fallback_box.text_frame
            ftf.word_wrap = True
            fp = ftf.paragraphs[0]
            fp.text = f"[Screenshot: {image_filename}]"
            fp.font.name = 'Arial'
            fp.font.size = Pt(18)
            fp.font.bold = True
            fp.font.color.rgb = text_muted

    slide_layout = prs.slide_layouts[6] # Blank

    # ────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Dark Theme)
    # ────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1, navy)
    
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "RIS Smart Recruitment & Roster System"
    p1.font.name = 'Arial'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = white
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Standardizing candidate intake, operational vacancy control, and structured Excel reporting"
    p1_sub.font.name = 'Arial'
    p1_sub.font.size = Pt(20)
    p1_sub.font.color.rgb = RGBColor(194, 205, 219)
    p1_sub.space_before = Pt(15)

    p1_foot = tf1.add_paragraph()
    p1_foot.text = "Executive Briefing  •  Solutions & System Walkthrough"
    p1_foot.font.name = 'Arial'
    p1_foot.font.size = Pt(12)
    p1_foot.font.color.rgb = text_muted
    p1_foot.space_before = Pt(80)

    # ────────────────────────────────────────────────────────
    # SLIDE 2: The Core Challenge: Legacies & Delimiter Sprawls
    # ────────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(slide_layout)
    set_background(slide2, light_gray)
    add_slide_title(slide2, "The Legacy Challenge: The Data & Roster Bottleneck")
    
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    lp1 = ltf.paragraphs[0]
    lp1.text = "THE RECRUITMENT BOTTLENECK"
    lp1.font.name = 'Arial'
    lp1.font.size = Pt(14)
    lp1.font.bold = True
    lp1.font.color.rgb = red
    
    lp2 = ltf.add_paragraph()
    lp2.text = "Data Scrawl & Manual Review"
    lp2.font.name = 'Arial'
    lp2.font.size = Pt(24)
    lp2.font.bold = True
    lp2.font.color.rgb = navy
    lp2.space_before = Pt(8)
    lp2.space_after = Pt(15)
    
    problems = [
        "Multiple degrees and research works clumped into single cells with commas, blocking automated analysis.",
        "Spreadsheets stretching infinitely to the right with redundant columns, causing horizontal scrolling fatigue.",
        "Inconsistent score formats (CGPAs vs. percentages) making comparison difficult.",
        "Manual resume downloads and alignment causing long hiring delays."
    ]
    for prob in problems:
        p = ltf.add_paragraph()
        p.text = "• " + prob
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = text_dark
        p.space_after = Pt(10)

    # Right Card: The Solution Summary
    right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp1 = rtf.paragraphs[0]
    rp1.text = "THE INTEGRATED PORTAL APPROACH"
    rp1.font.name = 'Arial'
    rp1.font.size = Pt(14)
    rp1.font.bold = True
    rp1.font.color.rgb = green
    
    rp2 = rtf.add_paragraph()
    rp2.text = "Structured Intake to Clean Output"
    rp2.font.name = 'Arial'
    rp2.font.size = Pt(24)
    rp2.font.bold = True
    rp2.font.color.rgb = navy
    rp2.space_before = Pt(8)
    rp2.space_after = Pt(15)
    
    solutions = [
        "Standardized Intake: Candidates input individual degrees, scores, and timelines dynamically at source.",
        "HR Control Centers: Instantly filter active roles, review metrics, and open dossier profile views.",
        "Stacked Roster Layouts: Candidate details are stacked vertically with merged basic details.",
        "Score Normalization: Auto-formats scores to clean percentages ('79%') or CGPA ('8 CGPA') on download."
    ]
    for sol in solutions:
        p = rtf.add_paragraph()
        p.text = "• " + sol
        p.font.name = 'Arial'
        p.font.size = Pt(14)
        p.font.color.rgb = text_dark
        p.space_after = Pt(10)

    # ────────────────────────────────────────────────────────
    # SLIDE 3: Candidate Experience: Job Listings
    # ────────────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(slide_layout)
    set_background(slide3, light_gray)
    add_text_and_image(
        slide=slide3,
        title="Candidate Interface: Structured Job Board",
        bullets=[
            ("Explore Opportunities", "Applicants access a clean web interface listing open opportunities at RIS."),
            ("Transparent Information", "Displays role titles, division tags (e.g. CMEC, DAKSHIN), locations, and submission deadlines."),
            ("Clear Scope", "Provides concise public descriptions to inform prospective applicants prior to starting an application."),
            ("Frictionless Flow", "Direct redirection to a structured application form to ensure candidate details are captured correctly.")
        ],
        image_filename="media__1779989599270.png"
    )

    # ────────────────────────────────────────────────────────
    # SLIDE 4: HR Job Dashboard
    # ────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(slide_layout)
    set_background(slide4, light_gray)
    add_text_and_image(
        slide=slide4,
        title="HR Dashboard: Central Roster Management",
        bullets=[
            ("Real-Time Tracking", "HR administrators get a comprehensive overview of active vacancies and applicant volumes."),
            ("Dashboard Statistics", "Top-level summary cards display Open Positions, Total Applicants (this year), and roles Closing Soon."),
            ("Status Columns", "Central table showing active vacancies, divisions, applicant numbers, posted dates, and deadlines."),
            ("Direct Operations", "Quick access links to view applicants, edit job listings, or archive roles.")
        ],
        image_filename="media__1779989675165.png"
    )

    # ────────────────────────────────────────────────────────
    # SLIDE 5: Job Creation Modal
    # ────────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(slide_layout)
    set_background(slide5, light_gray)
    add_text_and_image(
        slide=slide5,
        title="HR Operations: Standardized Vacancy Control",
        bullets=[
            ("Structured Input Forms", "Clean job creation modal with simplified fields to prevent layout clutter."),
            ("Compensation Settings", "Dedicated contract and compensation settings panel to manage role boundaries."),
            ("Predefined Pay Bands", "Set minimum and maximum pay expectations directly via dropdown selectors."),
            ("Experience Requirements", "Enforce minimum and maximum experience levels required, and set contract durations.")
        ],
        image_filename="media__1779989658124.png"
    )

    # ────────────────────────────────────────────────────────
    # SLIDE 6: System Analytics Overview
    # ────────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(slide_layout)
    set_background(slide6, light_gray)
    add_text_and_image(
        slide=slide6,
        title="HR Dashboard: System Talent Analytics",
        bullets=[
            ("Operational Metrics", "Presents visual candidate demographics across the system automatically."),
            ("Gender Diversity", "A donut chart maps the gender balance of the active applicant pool."),
            ("Top Talent Hubs", "Visualizes geographic distribution, showing top applicant states at a glance."),
            ("Hiring Pipeline Velocity", "Tracks total applications against shortlists and selections dynamically.")
        ],
        image_filename="media__1779989707679.png"
    )

    # ────────────────────────────────────────────────────────
    # SLIDE 7: Job-Specific Analytics & Filters
    # ────────────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(slide_layout)
    set_background(slide7, light_gray)
    add_text_and_image(
        slide=slide7,
        title="Roster Analytics & Multivariate Filters",
        bullets=[
            ("Vacancy Deep-Dive", "Displays demographic and seniority distributions tailored to the selected job posting."),
            ("Seniority Ratios", "Visualizes the distribution of Ph.D., Masters, and Bachelors degrees in the pool."),
            ("Advanced Filters", "Isolate candidate subsets instantly using biographical, academic, experience, and publication filters."),
            ("No Decision Lag", "Enables HR to search through hundreds of profiles in seconds without opening individual resumes.")
        ],
        image_filename="media__1779989728466.png"
    )

    # ────────────────────────────────────────────────────────
    # SLIDE 8: Candidate Dossier Views (Double Image: Side-by-Side)
    # ────────────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(slide_layout)
    set_background(slide8, light_gray)
    add_slide_title(slide8, "Candidate Dossier: 360-Degree Profile View")
    
    # Left Box Text
    t_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8))
    tf8 = t_box.text_frame
    tf8.word_wrap = True
    p8 = tf8.paragraphs[0]
    p8.text = "Review candidate profiles instantly in a modular, color-coded dossier modal without opening separate files."
    p8.font.name = 'Arial'
    p8.font.size = Pt(15)
    p8.font.color.rgb = text_muted
    
    # Embed Left Image (Karan Gupta top)
    brain_dir = "C:/Users/Viraal/.gemini/antigravity-ide/brain/b52771bd-afd4-41e0-baec-88c88ae23617"
    img_left = "media__1779989857660.png"
    img_left_path = os.path.join(brain_dir, img_left)
    if os.path.exists(img_left_path):
        slide8.shapes.add_picture(img_left_path, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.2))
        
    # Embed Right Image (Karan Gupta bottom)
    img_right = "media__1779989869625.png"
    img_right_path = os.path.join(brain_dir, img_right)
    if os.path.exists(img_right_path):
        slide8.shapes.add_picture(img_right_path, Inches(6.8), Inches(2.2), Inches(5.6), Inches(4.2))

    # Footer note on slide
    note_box = slide8.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5))
    ntf = note_box.text_frame
    np = ntf.paragraphs[0]
    np.text = "• Quick Profile Rollups  • Color-Coded Academic Levels  • Timeline Experience Views  • Clean Close Layout"
    np.font.name = 'Arial'
    np.font.size = Pt(13)
    np.font.bold = True
    np.font.color.rgb = indigo

    # ────────────────────────────────────────────────────────
    # SLIDE 9: Roster Export Options
    # ────────────────────────────────────────────────────────
    slide9 = prs.slides.add_slide(slide_layout)
    set_background(slide9, light_gray)
    add_text_and_image(
        slide=slide9,
        title="Roster Exports: Clean Reports Generation",
        bullets=[
            ("Flexible Formats", "Download candidate data tailored to the requirements of the review committee."),
            ("Detailed Grouped Excel", "Candidate qualifications are stacked vertically with basic details merged, keeping layouts clean."),
            ("Standardized Summary", "Generates condensed reports displaying degrees and latest roles inline (e.g. 'Degree (University)')."),
            ("Clean Score Normalization", "Automatically formats higher education scores ('79%' or '8 CGPA'), stripping trailing '.0' decimals.")
        ],
        image_filename="media__1779989923219.png"
    )

    # ────────────────────────────────────────────────────────
    # SLIDE 10: Strategic Benefits (Dark Background)
    # ────────────────────────────────────────────────────────
    slide10 = prs.slides.add_slide(slide_layout)
    set_background(slide10, navy)
    add_slide_title(slide10, "Smart Recruitment Portal: Operational Impact", white)
    
    sb_box = slide10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    sbtf = sb_box.text_frame
    sbtf.word_wrap = True
    
    benefits = [
        ("Time Saved in Review", "Managers review candidates directly from the portal, eliminating manual CV tracking and administrative overhead."),
        ("Standardized Candidate Data", "Dynamic form validations ensure applicant details are clean and complete on submission."),
        ("Advanced Roster Exports", "Generates presentation-ready spreadsheets with grouped stacked layouts, visual border dividers, and clean academic score normalization."),
        ("Board Room Ready", "Roster filters, dossier popups, and standardized exports are structured to be shared directly with senior board members and stakeholders.")
    ]
    for idx, (title, desc) in enumerate(benefits):
        p = sbtf.add_paragraph()
        p.text = f"{title}: "
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(194, 205, 219)
        p.space_before = Pt(18)
        
        run = p.add_run()
        run.text = desc
        run.font.name = 'Arial'
        run.font.size = Pt(15)
        run.font.bold = False
        run.font.color.rgb = RGBColor(226, 232, 240)

    # Save presentation
    output_dir = "c:/Users/Viraal/Desktop/HRForm/resources"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "HR_Portal_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation generated successfully with screenshots: {output_path}")

if __name__ == '__main__':
    create_presentation()
